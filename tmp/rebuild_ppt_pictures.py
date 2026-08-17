from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Iterable

from pptx import Presentation


KB_ROOT = Path(r"H:\xisixiang")
PPT_ROOT = Path(r"C:\Users\18105\Downloads\2025年秋教育部 第一章    新时代坚持和发展中国特色社会主义等13个文件")
WORK_ROOT = KB_ROOT / "tmp" / "ppt_picture_rebuild"
MANIFEST_PATH = WORK_ROOT / "selection_manifest.json"

CHAPTERS = [
    ("导论", "2025年秋季教育部导论习近平新时代中国特色社会主义思想概论概述课件.pptx", ["0001", "0002", "0003", "0004", "0005", "0006"], "intro"),
    ("第一章", "2025年秋教育部 第一章    新时代坚持和发展中国特色社会主义.pptx", [f"01{i:02d}" for i in range(1, 10)], "ch01"),
    ("第二章", "2025秋季教育部第二章 以中国式现代化全面推进中华民族伟大复兴.pptx", [f"02{i:02d}" for i in range(1, 11)], "ch02"),
    ("第三章", "2025年秋教育部第三章 坚持党的全面领导(1).pptx", [f"03{i:02d}" for i in range(1, 10)], "ch03"),
    ("第五章", "2025年秋教育部第五章 全面深化改革.pptx", [f"05{i:02d}" for i in range(1, 10)], "ch05"),
    ("第六章", "2025年秋教育部第六章推动高质量发展.pptx", [f"06{i:02d}" for i in range(1, 14)], "ch06"),
    ("第七章", "2025年秋教育部第七章    社会主义现代化建设的教育、科技、人才战略.pptx", [f"07{i:02d}" for i in range(1, 13)], "ch07"),
    ("第十章", "2025年秋教育部第十章 建设社会主义文化强国.pptx", [f"10{i:02d}" for i in range(1, 13)], "ch10"),
    ("第十二章", "2025年秋教育部第十二章  建设社会主义生态文明(1).pptx", [f"12{i:02d}" for i in range(1, 10)], "ch12"),
    ("第十三章", "2025年秋教育部第十三章 维护和塑造国家安全修改版(1).pptx", [f"13{i:02d}" for i in range(1, 10)], "ch13"),
    ("第十六章", "2025年秋教育部第十六章 中国特色大国外交和推动构建人类命运共同体.pptx", [f"16{i:02d}" for i in range(1, 12)], "ch16"),
    ("第十七章", "2025年秋教育部第十七章 全面从严治党.pptx", [f"17{i:02d}" for i in range(1, 15)], "ch17"),
]


def iter_shapes(shapes: Iterable) -> Iterable:
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from iter_shapes(shape.shapes)


def shape_text(shape) -> str:
    chunks: list[str] = []
    if getattr(shape, "has_text_frame", False):
        chunks.append(shape.text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            chunks.extend(cell.text for cell in row.cells)
    return "\n".join(x for x in chunks if x and x.strip())


def normalize(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def node_terms(node: dict) -> list[str]:
    values = [node.get("name", ""), *node.get("aliases", []), *node.get("keywords", []), *node.get("tags", [])]
    return sorted({normalize(str(v)) for v in values if len(normalize(str(v))) >= 2}, key=len, reverse=True)


def choose_node(text: str, nodes: list[dict]) -> tuple[dict | None, int, list[str]]:
    best: tuple[int, dict | None, list[str]] = (0, None, [])
    for node in nodes:
        terms = node_terms(node)
        hits = [term for term in terms if term in text]
        score = sum(3 if len(term) >= 5 else 2 for term in hits)
        if node.get("name") and node["name"] in text:
            score += 5
        if score > best[0]:
            best = (score, node, hits[:8])
    return best[1], best[0], best[2]


def main() -> None:
    WORK_ROOT.mkdir(parents=True, exist_ok=True)
    nodes_by_id = {
        path.stem.removeprefix("node_"): json.loads(path.read_text(encoding="utf-8"))
        for path in (KB_ROOT / "nodes").glob("node_*.json")
    }
    selections: list[dict] = []
    chapter_stats: list[dict] = []
    for chapter, filename, node_suffixes, slug in CHAPTERS:
        ppt_path = PPT_ROOT / filename
        candidate_nodes = [nodes_by_id[suffix] for suffix in node_suffixes if suffix in nodes_by_id]
        presentation = Presentation(str(ppt_path))
        selected = 0
        for slide_number, slide in enumerate(presentation.slides, start=1):
            text = normalize("\n".join(shape_text(shape) for shape in iter_shapes(slide.shapes)))
            if len(text) < 40:
                continue
            node, score, hits = choose_node(text, candidate_nodes)
            if node is None or score < 3:
                continue
            selected += 1
            selections.append({
                "chapter": chapter,
                "slug": slug,
                "ppt": str(ppt_path),
                "slide": slide_number,
                "node_id": node["id"],
                "node_name": node["name"],
                "node_abstract": node["abstract"],
                "matched_terms": hits,
                "text": text,
            })
        chapter_stats.append({"chapter": chapter, "ppt": str(ppt_path), "slides": len(presentation.slides), "selected": selected})
    MANIFEST_PATH.write_text(json.dumps({"chapters": chapter_stats, "selections": selections}, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({"chapters": chapter_stats, "selected_total": len(selections), "manifest": str(MANIFEST_PATH)}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
